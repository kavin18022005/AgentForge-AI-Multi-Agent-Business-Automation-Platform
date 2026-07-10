import os
import uuid
from loguru import logger
from sqlalchemy.ext.asyncio import AsyncSession
from app.models.report import Report

# ──────────────────────────────────────────────────────────────────────
# Report Exporters
# ──────────────────────────────────────────────────────────────────────

async def generate_export(report: Report, format: str, db: AsyncSession) -> str:
    """Generate a report document in the specified format and return the disk path."""
    project_id = str(report.project_id)
    report_id = str(report.id)
    
    # Define exports directory
    exports_dir = os.path.join("./uploads", "exports", project_id)
    os.makedirs(exports_dir, exist_ok=True)
    
    file_path = os.path.join(exports_dir, f"{report_id}.{format}")
    
    # If the file already exists, return its path
    if os.path.exists(file_path):
        return file_path
        
    try:
        if format == "md":
            await generate_md(report, file_path)
        elif format == "docx":
            await generate_docx(report, file_path)
        elif format == "pptx":
            await generate_pptx(report, file_path)
        elif format == "pdf":
            await generate_pdf(report, file_path)
        else:
            raise ValueError(f"Unsupported format: {format}")
            
        logger.info(f"Successfully generated {format} report export at {file_path}")
        
        # Save path in DB dynamically
        path_attr = f"{format}_path"
        if hasattr(report, path_attr):
            setattr(report, path_attr, file_path)
            await db.commit()
            
        return file_path
    except Exception as e:
        logger.error(f"Failed to generate {format} export: {e}")
        raise e


async def generate_md(report: Report, file_path: str):
    """Generate Markdown report file."""
    content = f"# {report.title}\n\n"
    content += f"**Report Type:** {report.report_type.replace('_', ' ').title()}\n"
    content += f"**Quality Score:** {report.quality_score}/10\n"
    content += f"**Date:** {report.created_at.strftime('%Y-%m-%d %H:%M:%S')}\n\n"
    content += "## Executive Summary\n\n"
    content += f"{report.executive_summary}\n\n"
    
    if report.content:
        content += "## Detailed Analysis\n\n"
        for section, data in report.content.items():
            if not data or section == "qa":
                continue
            content += f"### {section.replace('_', ' ').title()}\n\n"
            if isinstance(data, dict):
                for k, v in data.items():
                    content += f"**{k.replace('_', ' ').title()}:** {v}\n\n"
            elif isinstance(data, list):
                for item in data:
                    if isinstance(item, dict):
                        for k, v in item.items():
                            content += f"- **{k.replace('_', ' ').title()}:** {v}\n"
                        content += "\n"
                    else:
                        content += f"- {item}\n"
                content += "\n"
            else:
                content += f"{data}\n\n"
                
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(content)


async def generate_docx(report: Report, file_path: str):
    """Generate Word (DOCX) report file."""
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    
    doc = Document()
    
    # Document Title
    title = doc.add_heading(report.title, level=0)
    title.style.font.size = Pt(28)
    title.style.font.name = "Arial"
    
    # Metadata
    meta = doc.add_paragraph()
    meta.add_run(f"Report Type: ").bold = True
    meta.add_run(f"{report.report_type.replace('_', ' ').title()}\n")
    meta.add_run(f"Quality Score: ").bold = True
    meta.add_run(f"{report.quality_score}/10\n")
    meta.add_run(f"Word Count: ").bold = True
    meta.add_run(f"{report.word_count} words\n")
    meta.add_run(f"Date: ").bold = True
    meta.add_run(f"{report.created_at.strftime('%Y-%m-%d')}\n")
    
    doc.add_page_break()
    
    # Executive Summary Section
    h1 = doc.add_heading("Executive Summary", level=1)
    h1.style.font.color.rgb = RGBColor(99, 102, 241) # Brand Primary
    
    # Split paragraphs by newline to keep formatting
    for para in (report.executive_summary or "").split("\n\n"):
        if para.strip():
            doc.add_paragraph(para.strip())
            
    # Add detailed components if available
    if report.content:
        doc.add_page_break()
        h2 = doc.add_heading("Detailed Analysis", level=1)
        h2.style.font.color.rgb = RGBColor(99, 102, 241)
        
        for section, data in report.content.items():
            if not data or section == "qa":
                continue
                
            doc.add_heading(section.replace('_', ' ').title(), level=2)
            
            if isinstance(data, dict):
                for k, v in data.items():
                    p = doc.add_paragraph()
                    p.add_run(f"{k.replace('_', ' ').title()}: ").bold = True
                    p.add_run(str(v))
            elif isinstance(data, list):
                for item in data:
                    if isinstance(item, dict):
                        p = doc.add_paragraph(style="List Bullet")
                        for idx, (k, v) in enumerate(item.items()):
                            sep = " | " if idx > 0 else ""
                            p.add_run(f"{sep}{k.replace('_', ' ').title()}: ").bold = True
                            p.add_run(str(v))
                    else:
                        doc.add_paragraph(str(item), style="List Bullet")
            else:
                doc.add_paragraph(str(data))
                
    doc.save(file_path)


async def generate_pptx(report: Report, file_path: str):
    """Generate PowerPoint (PPTX) report file."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    prs = Presentation()
    
    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = report.title
    slide.placeholders[1].text = f"AgentForge AI Intelligence Report\nType: {report.report_type.replace('_', ' ').title()}\nScore: {report.quality_score}/10"
    
    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    
    # Cut summary to fit slide
    summary_text = report.executive_summary or ""
    paragraphs = [p.strip() for p in summary_text.split("\n\n") if p.strip()]
    for idx, p in enumerate(paragraphs[:3]):
        if idx == 0:
            tf.paragraphs[0].text = p
        else:
            p_node = tf.add_paragraph()
            p_node.text = p
            p_node.space_before = Pt(10)
            
    # Add detailed modules as slides
    if report.content:
        for section, data in report.content.items():
            if not data or section == "qa" or section == "manager":
                continue
                
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section.replace('_', ' ').title()
            tf = slide.placeholders[1].text_frame
            tf.word_wrap = True
            
            if isinstance(data, dict):
                items = list(data.items())[:5]
                for idx, (k, v) in enumerate(items):
                    val_str = str(v)
                    if len(val_str) > 100:
                        val_str = val_str[:97] + "..."
                    line = f"• {k.replace('_', ' ').title()}: {val_str}"
                    if idx == 0:
                        tf.paragraphs[0].text = line
                    else:
                        p_node = tf.add_paragraph()
                        p_node.text = line
                        p_node.space_before = Pt(6)
            elif isinstance(data, list):
                items = data[:5]
                for idx, item in enumerate(items):
                    if isinstance(item, dict):
                        # Combine dictionary values into single line
                        line_parts = [f"{k.replace('_',' ').title()}: {v}" for k, v in item.items() if v]
                        line = "• " + " | ".join(line_parts[:3])
                    else:
                        line = f"• {str(item)}"
                        
                    if len(line) > 120:
                        line = line[:117] + "..."
                        
                    if idx == 0:
                        tf.paragraphs[0].text = line
                    else:
                        p_node = tf.add_paragraph()
                        p_node.text = line
                        p_node.space_before = Pt(6)
            else:
                desc = str(data)
                if len(desc) > 400:
                    desc = desc[:397] + "..."
                tf.paragraphs[0].text = desc

    prs.save(file_path)


async def generate_pdf(report: Report, file_path: str):
    """Generate PDF report file using ReportLab platypus flowables."""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    
    doc = SimpleDocTemplate(
        file_path,
        pagesize=letter,
        rightMargin=54,
        leftMargin=54,
        topMargin=54,
        bottomMargin=54
    )
    
    styles = getSampleStyleSheet()
    
    # Custom Brand Colors & Styles
    primary_color = colors.HexColor("#6366f1")   # indigo-500
    text_color = colors.HexColor("#1e293b")      # slate-800
    secondary_color = colors.HexColor("#475569") # slate-600
    
    title_style = ParagraphStyle(
        "PDFTitle",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=26,
        leading=32,
        textColor=primary_color,
        spaceAfter=20,
        alignment=0 # Left-align
    )
    
    h1_style = ParagraphStyle(
        "PDFH1",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=18,
        leading=22,
        textColor=primary_color,
        spaceBefore=18,
        spaceAfter=10,
        keepWithNext=True
    )
    
    h2_style = ParagraphStyle(
        "PDFH2",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=13,
        leading=16,
        textColor=secondary_color,
        spaceBefore=12,
        spaceAfter=6,
        keepWithNext=True
    )
    
    body_style = ParagraphStyle(
        "PDFBody",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=10.5,
        leading=15,
        textColor=text_color,
        spaceAfter=8
    )
    
    meta_style = ParagraphStyle(
        "PDFMeta",
        parent=styles["Normal"],
        fontName="Helvetica-Oblique",
        fontSize=9.5,
        leading=14,
        textColor=secondary_color,
        spaceAfter=15
    )
    
    story = []
    
    # --- Title Page / Block ---
    story.append(Paragraph(report.title, title_style))
    story.append(Spacer(1, 10))
    
    meta_text = (
        f"<b>Report Type:</b> {report.report_type.replace('_', ' ').title()}<br/>"
        f"<b>Quality Rating:</b> {report.quality_score}/10<br/>"
        f"<b>Scope:</b> {report.word_count} Words<br/>"
        f"<b>Generated on:</b> {report.created_at.strftime('%B %d, %Y')}"
    )
    story.append(Paragraph(meta_text, meta_style))
    story.append(Spacer(1, 20))
    
    # --- Executive Summary ---
    story.append(Paragraph("Executive Summary", h1_style))
    for paragraph_text in (report.executive_summary or "").split("\n\n"):
        if paragraph_text.strip():
            story.append(Paragraph(paragraph_text.strip(), body_style))
            
    # --- Details ---
    if report.content:
        story.append(PageBreak())
        story.append(Paragraph("Detailed Analysis Findings", h1_style))
        story.append(Spacer(1, 10))
        
        for section, data in report.content.items():
            if not data or section in ["qa", "manager"]:
                continue
                
            story.append(Paragraph(section.replace('_', ' ').title(), h2_style))
            
            if isinstance(data, dict):
                for k, v in data.items():
                    line = f"<b>{k.replace('_', ' ').title()}:</b> {v}"
                    story.append(Paragraph(line, body_style))
            elif isinstance(data, list):
                for item in data:
                    if isinstance(item, dict):
                        parts = [f"<b>{k.replace('_',' ').title()}:</b> {v}" for k, v in item.items() if v]
                        line = "• " + " | ".join(parts)
                        story.append(Paragraph(line, body_style))
                    else:
                        story.append(Paragraph(f"• {str(item)}", body_style))
            else:
                story.append(Paragraph(str(data), body_style))
                
            story.append(Spacer(1, 10))
            
    # Build Document flowable story
    doc.build(story)
